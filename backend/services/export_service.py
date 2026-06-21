"""Export Service — generates PDF, PPTX, Excel reports from dashboard state."""
import io
import json
import os
from datetime import datetime
from typing import Any, Dict, List
from loguru import logger


class ExportService:

    def export_excel(self, state: Dict[str, Any]) -> bytes:
        import xlsxwriter
        output = io.BytesIO()
        wb = xlsxwriter.Workbook(output, {"in_memory": True})

        # Formats
        title_fmt = wb.add_format({"bold": True, "font_size": 14, "bg_color": "#1E293B", "font_color": "#FFFFFF", "border": 1})
        header_fmt = wb.add_format({"bold": True, "bg_color": "#334155", "font_color": "#FFFFFF", "border": 1})
        data_fmt = wb.add_format({"border": 1, "font_color": "#1E293B"})
        number_fmt = wb.add_format({"border": 1, "num_format": "#,##0.00"})
        pct_fmt = wb.add_format({"border": 1, "num_format": "0.00%"})

        # Sheet 1: Executive Summary
        ws_exec = wb.add_worksheet("Executive Summary")
        ws_exec.set_column("A:A", 30)
        ws_exec.set_column("B:B", 50)

        ws_exec.write("A1", "PowerBI Genius AI — Executive Report", title_fmt)
        ws_exec.write("A2", f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}", data_fmt)
        ws_exec.write("A3", f"Domain: {state.get('domain', 'N/A').upper()}", data_fmt)

        summary = state.get("executive_summary", "No summary available.")
        ws_exec.write("A5", "Executive Summary", header_fmt)
        ws_exec.write_string("A6", summary, data_fmt)
        ws_exec.set_row(5, 200)

        # Sheet 2: KPIs
        ws_kpi = wb.add_worksheet("KPIs")
        ws_kpi.write_row("A1", ["KPI Name", "Value", "Category", "Description"], header_fmt)
        ws_kpi.set_column("A:A", 25)
        ws_kpi.set_column("B:B", 20)
        ws_kpi.set_column("C:C", 20)
        ws_kpi.set_column("D:D", 40)

        for row, kpi in enumerate(state.get("kpis", []), start=1):
            ws_kpi.write(row, 0, kpi.get("display_name", ""), data_fmt)
            ws_kpi.write(row, 1, kpi.get("formatted_value", ""), data_fmt)
            ws_kpi.write(row, 2, kpi.get("category", ""), data_fmt)
            ws_kpi.write(row, 3, kpi.get("description", ""), data_fmt)

        # Sheet 3: Insights
        ws_ins = wb.add_worksheet("AI Insights")
        ws_ins.write_row("A1", ["Title", "Category", "Impact", "Description", "Recommendation"], header_fmt)
        ws_ins.set_column("A:E", 35)

        for row, ins in enumerate(state.get("insights", []), start=1):
            ws_ins.write(row, 0, ins.get("title", ""), data_fmt)
            ws_ins.write(row, 1, ins.get("category", ""), data_fmt)
            ws_ins.write(row, 2, ins.get("impact", ""), data_fmt)
            ws_ins.write(row, 3, ins.get("description", ""), data_fmt)
            ws_ins.write(row, 4, ins.get("recommendation", ""), data_fmt)

        # Sheet 4: Cleaned Data
        cleaned = state.get("cleaned_data", {})
        if cleaned and cleaned.get("data"):
            ws_data = wb.add_worksheet("Cleaned Data")
            columns = cleaned.get("columns", [])
            ws_data.write_row(0, 0, columns, header_fmt)
            for row, record in enumerate(cleaned["data"][:10000], start=1):
                for col_idx, col in enumerate(columns):
                    val = record.get(col, "")
                    ws_data.write(row, col_idx, val if val is not None else "", data_fmt)

        # Sheet 5: DAX Measures
        if state.get("data_model"):
            ws_dax = wb.add_worksheet("DAX Measures")
            ws_dax.set_column("A:A", 30)
            ws_dax.set_column("B:B", 80)
            ws_dax.write_row("A1", ["Measure Name", "DAX Expression", "Format"], header_fmt)
            for row, measure in enumerate(state["data_model"].get("dax_measures", []), start=1):
                ws_dax.write(row, 0, measure.get("name", ""), data_fmt)
                ws_dax.write(row, 1, measure.get("expression", ""), data_fmt)
                ws_dax.write(row, 2, measure.get("format_string", ""), data_fmt)

        wb.close()
        return output.getvalue()

    def export_pdf(self, state: Dict[str, Any]) -> bytes:
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
        from reportlab.lib.enums import TA_CENTER, TA_LEFT

        output = io.BytesIO()
        doc = SimpleDocTemplate(output, pagesize=A4, rightMargin=2*cm, leftMargin=2*cm, topMargin=2*cm, bottomMargin=2*cm)

        styles = getSampleStyleSheet()
        dark_bg = colors.HexColor("#1E293B")
        accent = colors.HexColor("#6366F1")
        text_color = colors.HexColor("#1E293B")

        title_style = ParagraphStyle("Title", fontSize=24, textColor=accent, spaceAfter=6, alignment=TA_CENTER, fontName="Helvetica-Bold")
        sub_style = ParagraphStyle("Sub", fontSize=12, textColor=colors.grey, spaceAfter=20, alignment=TA_CENTER)
        section_style = ParagraphStyle("Section", fontSize=14, textColor=accent, spaceBefore=16, spaceAfter=8, fontName="Helvetica-Bold")
        body_style = ParagraphStyle("Body", fontSize=10, textColor=text_color, spaceAfter=6, leading=14)

        story = []

        # Title
        story.append(Spacer(1, 1*cm))
        domain = state.get("domain", "Business").upper()
        story.append(Paragraph(f"PowerBI Genius AI", title_style))
        story.append(Paragraph(f"{domain} Intelligence Dashboard Report", sub_style))
        story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}", sub_style))
        story.append(HRFlowable(width="100%", thickness=2, color=accent))
        story.append(Spacer(1, 0.5*cm))

        # Executive Summary
        story.append(Paragraph("Executive Summary", section_style))
        summary = state.get("executive_summary", "No summary generated.")
        story.append(Paragraph(summary, body_style))
        story.append(Spacer(1, 0.5*cm))

        # KPIs Table
        story.append(Paragraph("Key Performance Indicators", section_style))
        kpis = state.get("kpis", [])
        if kpis:
            kpi_data = [["KPI", "Value", "Category"]]
            for kpi in kpis[:8]:
                kpi_data.append([kpi.get("display_name", ""), kpi.get("formatted_value", "N/A"), kpi.get("category", "")])

            kpi_table = Table(kpi_data, colWidths=[7*cm, 4*cm, 5*cm])
            kpi_table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), dark_bg),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 9),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#F8FAFC"), colors.white]),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CBD5E1")),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]))
            story.append(kpi_table)

        # Insights
        story.append(Spacer(1, 0.5*cm))
        story.append(Paragraph("AI-Generated Insights", section_style))
        for ins in state.get("insights", [])[:6]:
            story.append(Paragraph(f"<b>{ins.get('title', '')}</b>", body_style))
            story.append(Paragraph(ins.get("description", ""), body_style))
            story.append(Paragraph(f"<i>Recommendation: {ins.get('recommendation', '')}</i>", body_style))
            story.append(Spacer(1, 0.3*cm))

        # Narrative
        narrative = state.get("narrative", "")
        if narrative:
            story.append(Paragraph("Strategic Narrative", section_style))
            story.append(Paragraph(narrative, body_style))

        doc.build(story)
        return output.getvalue()

    def export_pptx(self, state: Dict[str, Any]) -> bytes:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        DARK = RGBColor(0x0F, 0x17, 0x2A)
        SURFACE = RGBColor(0x1E, 0x29, 0x3B)
        ACCENT = RGBColor(0x63, 0x66, 0xF1)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        GREY = RGBColor(0x94, 0xA3, 0xB8)

        blank_layout = prs.slide_layouts[6]

        def add_text_box(slide, text, left, top, width, height, font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
            tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
            return tb

        def set_bg(slide, color: RGBColor):
            from pptx.oxml.ns import qn
            from lxml import etree
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = color

        # Slide 1: Title
        slide1 = prs.slides.add_slide(blank_layout)
        set_bg(slide1, DARK)
        domain = state.get("domain", "Business").upper()
        add_text_box(slide1, "PowerBI Genius AI", 1, 1.5, 11, 1.2, font_size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text_box(slide1, f"{domain} Intelligence Dashboard", 1, 2.9, 11, 0.8, font_size=22, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide1, f"Generated: {datetime.now().strftime('%B %d, %Y')}", 1, 3.9, 11, 0.5, font_size=14, color=GREY, align=PP_ALIGN.CENTER)
        add_text_box(slide1, "AI-Powered | McKinsey-Grade Analytics", 1, 5.5, 11, 0.5, font_size=12, color=GREY, align=PP_ALIGN.CENTER)

        # Slide 2: KPIs
        slide2 = prs.slides.add_slide(blank_layout)
        set_bg(slide2, DARK)
        add_text_box(slide2, "Key Performance Indicators", 0.5, 0.3, 12, 0.7, font_size=24, bold=True, color=ACCENT)

        kpis = state.get("kpis", [])[:6]
        positions = [(0.3, 1.2), (4.5, 1.2), (8.7, 1.2), (0.3, 3.8), (4.5, 3.8), (8.7, 3.8)]
        for kpi, (x, y) in zip(kpis, positions):
            tb = slide2.shapes.add_textbox(Inches(x), Inches(y), Inches(4), Inches(2))
            tf = tb.text_frame
            tf.word_wrap = True
            # Value
            p1 = tf.paragraphs[0]
            r1 = p1.add_run()
            r1.text = kpi.get("formatted_value", "N/A")
            r1.font.size = Pt(28)
            r1.font.bold = True
            r1.font.color.rgb = ACCENT
            # Name
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = kpi.get("display_name", "")
            r2.font.size = Pt(11)
            r2.font.color.rgb = GREY

        # Slide 3: Executive Summary
        slide3 = prs.slides.add_slide(blank_layout)
        set_bg(slide3, DARK)
        add_text_box(slide3, "Executive Summary", 0.5, 0.3, 12, 0.7, font_size=24, bold=True, color=ACCENT)
        summary = state.get("executive_summary", "")[:1200]
        add_text_box(slide3, summary, 0.5, 1.2, 12.3, 5.5, font_size=11, color=WHITE)

        # Slide 4: AI Insights
        slide4 = prs.slides.add_slide(blank_layout)
        set_bg(slide4, DARK)
        add_text_box(slide4, "AI-Generated Insights", 0.5, 0.3, 12, 0.7, font_size=24, bold=True, color=ACCENT)
        y_pos = 1.2
        for ins in state.get("insights", [])[:4]:
            title = f"• {ins.get('title', '')}"
            add_text_box(slide4, title, 0.5, y_pos, 12, 0.4, font_size=12, bold=True, color=WHITE)
            desc = ins.get("description", "")[:200]
            add_text_box(slide4, desc, 0.8, y_pos + 0.4, 12, 0.4, font_size=10, color=GREY)
            y_pos += 1.0

        # Slide 5: Recommendations
        slide5 = prs.slides.add_slide(blank_layout)
        set_bg(slide5, DARK)
        add_text_box(slide5, "Strategic Recommendations", 0.5, 0.3, 12, 0.7, font_size=24, bold=True, color=ACCENT)
        y_pos = 1.2
        for ins in state.get("insights", [])[:5]:
            rec = f"→ {ins.get('recommendation', '')}"
            add_text_box(slide5, rec, 0.5, y_pos, 12, 0.5, font_size=11, color=WHITE)
            y_pos += 0.9

        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()

    def generate_pbix_template(self, state: Dict[str, Any]) -> Dict[str, Any]:
        """Returns a Power BI template manifest that can be imported."""
        dashboard_spec = state.get("dashboard_spec", {})
        data_model = state.get("data_model", {})
        kpis = state.get("kpis", [])

        return {
            "version": "1.0",
            "dashboard_id": dashboard_spec.get("dashboard_id", ""),
            "title": dashboard_spec.get("title", "AI Dashboard"),
            "domain": state.get("domain", "unknown"),
            "pages": dashboard_spec.get("pages", []),
            "data_model": data_model,
            "dax_measures": [k.get("dax_measure", "") for k in kpis if k.get("dax_measure")],
            "theme": state.get("theme", {}),
            "color_palette": dashboard_spec.get("color_palette", []),
            "instructions": [
                "1. Open Power BI Desktop",
                "2. Import the cleaned data file (CSV/Excel)",
                "3. Apply the data model relationships defined in 'data_model'",
                "4. Create DAX measures from the 'dax_measures' section",
                "5. Build each page following the 'pages' specification",
                "6. Apply the theme colors from 'color_palette'",
            ],
        }
